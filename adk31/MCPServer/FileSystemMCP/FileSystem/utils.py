"""
Shared Utilities for FileSystem MCP Package
===========================================
Centralizes path safety, format detection, read/write helpers, library checks,
and constants. Imported by other modules.

BASE_DIR is now sourced exclusively from config.py (which reads .env).
IGNORED_NAMES is now sourced exclusively from config.py (which reads .env).
"""

import base64
import csv
import json
import logging
import os
import sys
import platform
import subprocess
from datetime import datetime
from typing import Any, Dict, List, Optional


# =============================================================================
# BASE DIR — imported from config (no direct .env access here)
# =============================================================================
from .config import BASE_DIR, IGNORED_NAMES  # noqa: E402  (config.py sits at project root)

os.makedirs(BASE_DIR, exist_ok=True)

# =============================================================================
# LOGGING — log/{date}_{time}.log, created when the module is first imported
# =============================================================================
_LOG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "log")
os.makedirs(_LOG_DIR, exist_ok=True)

_log_filename = datetime.now().strftime("log_%Y%m%d_%H%M%S") + ".log"
_log_filepath = os.path.join(_LOG_DIR, _log_filename)

logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s  [%(module)s]  %(levelname)-8s  %(name)s  %(message)s",
    handlers=[
        logging.FileHandler(_log_filepath, encoding="utf-8"),
        logging.StreamHandler(sys.stderr),
    ],
)

logger = logging.getLogger("file_system_server")
logger.info(f"FileSystem MCP Server started. Base directory: {BASE_DIR}, Ignored Names: {IGNORED_NAMES}")
logger.info(f"Log file: {_log_filepath}")

# =============================================================================
# Optional Library Detection
# =============================================================================
HAS_DOCX = False
try:
    import docx as _docx_mod
    HAS_DOCX = True
except ImportError:
    logger.warning("python-docx not installed – .docx support disabled. (pip install python-docx)")

HAS_XLSX = False
try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    logger.warning("openpyxl not installed – .xlsx support disabled. (pip install openpyxl)")

HAS_PDF_READ = False
try:
    import pdfplumber
    HAS_PDF_READ = True
except ImportError:
    logger.warning("pdfplumber not installed – PDF read disabled. (pip install pdfplumber)")

HAS_PDF_WRITE = False
try:
    from reportlab.lib.pagesizes import letter as rl_letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph as RLParagraph
    from reportlab.lib.styles import getSampleStyleSheet as rl_styles
    HAS_PDF_WRITE = True
except ImportError:
    logger.warning("reportlab not installed – PDF write disabled. (pip install reportlab)")

HAS_PPTX = False
try:
    from pptx import Presentation as _PptxPresentation
    HAS_PPTX = True
except ImportError:
    logger.warning("python-pptx not installed – .pptx support disabled. (pip install python-pptx)")

# =============================================================================
# Path Helpers
# =============================================================================
def safe_path(path: str) -> str:
    """Resolve path relative to BASE_DIR and block directory traversal."""
    stripped = path.lstrip("/\\")
    abs_path = os.path.abspath(os.path.join(BASE_DIR, stripped))
    if not (abs_path == BASE_DIR or abs_path.startswith(BASE_DIR + os.sep)):
        raise ValueError(f"Access denied: '{path}' resolves outside the sandbox directory.")
    return abs_path

def get_ext(file_path: str) -> str:
    return os.path.splitext(file_path)[1].lower()

def _ensure_parent(file_path: str) -> None:
    parent = os.path.dirname(file_path)
    if parent:
        os.makedirs(parent, exist_ok=True)

# =============================================================================
# Format-aware READ Helpers
# =============================================================================
def _read_text(fp: str) -> str:
    with open(fp, "r", encoding="utf-8", errors="ignore") as fh:
        return fh.read()

def _read_json(fp: str) -> str:
    with open(fp, "r", encoding="utf-8") as fh:
        data = json.load(fh)
    return json.dumps(data, indent=2, ensure_ascii=False)

def _read_csv(fp: str) -> str:
    with open(fp, "r", encoding="utf-8", errors="ignore", newline="") as fh:
        reader = csv.DictReader(fh)
        rows = list(reader)
    return json.dumps(rows, indent=2, ensure_ascii=False)

def _read_docx(fp: str) -> str:
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    doc = _docx_mod.Document(fp)
    paragraphs = [p.text for p in doc.paragraphs]
    tables = []
    for tbl in doc.tables:
        for row in tbl.rows:
            tables.append([cell.text for cell in row.cells])
    result = "\n".join(paragraphs)
    if tables:
        result += "\n\n[Tables]\n" + json.dumps(tables, indent=2)
    return result

def _read_xlsx(fp: str) -> str:
    if not HAS_XLSX:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
    import openpyxl
    wb = openpyxl.load_workbook(fp, data_only=True)
    output: Dict[str, list] = {}
    for name in wb.sheetnames:
        ws = wb[name]
        output[name] = [list(row) for row in ws.iter_rows(values_only=True)]
    return json.dumps(output, indent=2, default=str)

def _read_pdf(fp: str) -> str:
    if not HAS_PDF_READ:
        raise RuntimeError("pdfplumber not installed. Run: pip install pdfplumber")
    import pdfplumber
    pages: List[str] = []
    with pdfplumber.open(fp) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            tables = page.extract_tables() or []
            chunk = f"--- Page {i} ---\n{text}"
            if tables:
                chunk += "\n[Tables]\n" + json.dumps(tables, indent=2)
            pages.append(chunk)
    return "\n\n".join(pages)

def _read_pptx(fp: str) -> str:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = _PptxPresentation(fp)
    slides: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        slides.append(f"--- Slide {i} ---\n" + "\n".join(lines))
    return "\n\n".join(slides)

def _read_binary_b64(fp: str) -> str:
    with open(fp, "rb") as fh:
        return base64.b64encode(fh.read()).decode("utf-8")

# =============================================================================
# Format-aware WRITE Helpers
# =============================================================================
def _write_text(fp: str, content: str) -> None:
    _ensure_parent(fp)
    with open(fp, "w", encoding="utf-8") as fh:
        fh.write(content)

def _write_json(fp: str, content: str) -> None:
    parsed = json.loads(content)
    _ensure_parent(fp)
    with open(fp, "w", encoding="utf-8") as fh:
        json.dump(parsed, fh, indent=2, ensure_ascii=False)

def _write_csv(fp: str, content: str) -> None:
    rows = json.loads(content)
    _ensure_parent(fp)
    with open(fp, "w", encoding="utf-8", newline="") as fh:
        if rows and isinstance(rows[0], dict):
            writer = csv.DictWriter(fh, fieldnames=list(rows[0].keys()))
            writer.writeheader()
            writer.writerows(rows)
        else:
            csv.writer(fh).writerows(rows)

def _write_docx(fp: str, content: str) -> None:
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    _ensure_parent(fp)
    lines = content.splitlines()
    if os.path.exists(fp):
        doc = _docx_mod.Document(fp)
        for i, para in enumerate(doc.paragraphs):
            para.clear()
            if i < len(lines):
                para.add_run(lines[i])
        for line in lines[len(doc.paragraphs):]:
            doc.add_paragraph(line)
    else:
        doc = _docx_mod.Document()
        for line in lines:
            doc.add_paragraph(line)
    doc.save(fp)

def _write_xlsx(fp: str, content: str) -> None:
    if not HAS_XLSX:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
    import openpyxl
    data = json.loads(content)
    if isinstance(data, list):
        data = {"Sheet1": data}
    _ensure_parent(fp)
    wb = openpyxl.load_workbook(fp) if os.path.exists(fp) else openpyxl.Workbook()
    if not os.path.exists(fp) and "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    for sheet_name, rows in data.items():
        if sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            ws.delete_rows(1, ws.max_row)
        else:
            ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    wb.save(fp)

def _write_pdf(fp: str, content: str) -> None:
    if not HAS_PDF_WRITE:
        raise RuntimeError("reportlab not installed. Run: pip install reportlab")
    _ensure_parent(fp)
    styles = rl_styles()
    doc = SimpleDocTemplate(fp, pagesize=rl_letter)
    story = [
        RLParagraph(line if line.strip() else "&nbsp;", styles["Normal"])
        for line in content.splitlines()
    ]
    doc.build(story)

def _write_pptx(fp: str, content: str) -> None:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    _ensure_parent(fp)
    try:
        slides_data = json.loads(content)
        if not isinstance(slides_data, list):
            raise ValueError
    except (json.JSONDecodeError, ValueError):
        import re
        parts = re.split(r"---\s*Slide\s*\d+\s*---", content)
        slides_data = [{"title": "", "body": p.strip()} for p in parts if p.strip()]
    prs = _PptxPresentation(fp) if os.path.exists(fp) else _PptxPresentation()
    layout = prs.slide_layouts[1]
    for info in slides_data:
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title and isinstance(info, dict):
            slide.shapes.title.text = info.get("title", "")
        body = slide.placeholders[1]
        if body:
            body.text = info.get("body", "") if isinstance(info, dict) else str(info)
    prs.save(fp)

def _write_binary_b64(fp: str, content: str) -> None:
    _ensure_parent(fp)
    with open(fp, "wb") as fh:
        fh.write(base64.b64decode(content))

# =============================================================================
# Extension Dispatch Tables and Constants
# =============================================================================
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst",
    ".html", ".htm", ".xml", ".svg",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs",
    ".sh", ".bash", ".zsh", ".bat", ".ps1",
    ".sql", ".psql", ".r", ".rb", ".php", ".swift", ".kt", ".dart",
    ".log", ".gitignore", ".dockerignore", ".editorconfig",
}

READ_DISPATCH = {
    ".json": _read_json,
    ".csv":  _read_csv,
    ".tsv":  _read_csv,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".xls":  _read_xlsx,
    ".pdf":  _read_pdf,
    ".pptx": _read_pptx,
}

WRITE_DISPATCH = {
    ".json": _write_json,
    ".csv":  _write_csv,
    ".tsv":  _write_csv,
    ".docx": _write_docx,
    ".xlsx": _write_xlsx,
    ".xls":  _write_xlsx,
    ".pdf":  _write_pdf,
    ".pptx": _write_pptx,
}