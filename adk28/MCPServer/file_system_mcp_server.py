"""
FileSystem MCP Server  v0.4.0
==============================
Full-featured file-system MCP server built on Google ADK + MCP STDIO.

HOW TO CONFIGURE:
  Set BASE_DIR below to the folder you want to work in.
  No command-line arguments needed.

CAPABILITIES
------------
File operations (any path, including sub-directories):
  • read_file        – read any file type (text, JSON, CSV, DOCX, XLSX, PDF, PPTX, binary)
  • write_file       – create or overwrite any file type
  • edit_file        – find-and-replace inside text-based files
  • append_file      – append text to a file
  • copy_file        – copy a file to another location
  • move_file        – move / rename a file
  • delete_file      – delete a single file
  • clear_file       – empty a text file without deleting it
  • file_info        – metadata (size, timestamps, format) for any file or directory

Directory operations:
  • create_directory  – create a folder (including nested sub-folders in one call)
  • delete_directory  – delete a folder (optionally recursive)
  • rename_directory  – rename / move a directory
  • list_files        – list files in a directory (non-recursive)
  • list_directories  – list immediate sub-directories
  • list_tree         – recursive tree view of a directory

Utility:
  • supported_formats – show all supported formats and library availability

Supported file formats
----------------------
Built-in : .txt .md .html .xml .yaml .yml .toml .ini .cfg .conf .log .env
           .py .js .ts .jsx .tsx .java .c .cpp .h .cs .go .rs
           .sh .bash .bat .ps1 .sql .psql .r .rb .php .swift .kt .dart
           .json  .csv  .tsv
Optional : .docx  (python-docx)
           .xlsx / .xls  (openpyxl)
           .pdf  read (pdfplumber)  write (reportlab)
           .pptx  (python-pptx)
Fallback : any unrecognised binary → base64 round-trip

INSTALL OPTIONAL DEPS
---------------------
pip install python-docx openpyxl pdfplumber reportlab python-pptx
"""

import asyncio
import base64
import csv
import json
import logging
import os
import shutil
import sys
from typing import Any, Dict, List

# --- MCP Server Imports ---
from mcp import types as mcp_types
from mcp.server.lowlevel import Server, NotificationOptions
from mcp.server.models import InitializationOptions
import mcp.server.stdio

# --- ADK Imports ---
from google.adk.tools.function_tool import FunctionTool
from google.adk.tools.mcp_tool.conversion_utils import adk_to_mcp_tool_type

# =============================================================================
# SET YOUR WORKING DIRECTORY HERE
# =============================================================================
BASE_DIR = r"D:\Agent-Development-Kit\adk23\Vikas"
# =============================================================================

BASE_DIR = os.path.abspath(BASE_DIR)
os.makedirs(BASE_DIR, exist_ok=True)

# =============================================================================
# Logging
# =============================================================================
logging.basicConfig(level=logging.INFO, stream=sys.stderr)
logger = logging.getLogger("file_system_server")
logger.info(f"FileSystem MCP Server started. Base directory: {BASE_DIR}")

# =============================================================================
# MCP App
# =============================================================================
app = Server("FileSystem-ADK-MCP")

# =============================================================================
# Optional library detection (server still starts if any are missing)
# =============================================================================
try:
    import docx as _docx_mod
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.warning("python-docx not installed – .docx support disabled. (pip install python-docx)")

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False
    logger.warning("openpyxl not installed – .xlsx support disabled. (pip install openpyxl)")

try:
    import pdfplumber
    HAS_PDF_READ = True
except ImportError:
    HAS_PDF_READ = False
    logger.warning("pdfplumber not installed – PDF read disabled. (pip install pdfplumber)")

try:
    from reportlab.lib.pagesizes import letter as rl_letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph as RLParagraph
    from reportlab.lib.styles import getSampleStyleSheet as rl_styles
    HAS_PDF_WRITE = True
except ImportError:
    HAS_PDF_WRITE = False
    logger.warning("reportlab not installed – PDF write disabled. (pip install reportlab)")

try:
    from pptx import Presentation as _PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed – .pptx support disabled. (pip install python-pptx)")

# =============================================================================
# Path Helpers
# =============================================================================

def safe_path(path: str) -> str:
    """Resolve path relative to BASE_DIR and block directory traversal."""
    stripped = path.lstrip("/\\")
    abs_path = os.path.abspath(os.path.join(BASE_DIR, stripped))
    if not (abs_path == BASE_DIR or abs_path.startswith(BASE_DIR + os.sep)):
        raise ValueError(f"Access denied: '{path}' resolves outside the sandbox directory.")
    return abs_path


def get_ext(file_path: str) -> str:
    return os.path.splitext(file_path)[1].lower()


def _ensure_parent(file_path: str) -> None:
    parent = os.path.dirname(file_path)
    if parent:
        os.makedirs(parent, exist_ok=True)


# =============================================================================
# Format-aware READ Helpers
# =============================================================================

def _read_text(fp: str) -> str:
    with open(fp, "r", encoding="utf-8", errors="ignore") as fh:
        return fh.read()


def _read_json(fp: str) -> str:
    with open(fp, "r", encoding="utf-8") as fh:
        data = json.load(fh)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _read_csv(fp: str) -> str:
    with open(fp, "r", encoding="utf-8", errors="ignore", newline="") as fh:
        reader = csv.DictReader(fh)
        rows = list(reader)
    return json.dumps(rows, indent=2, ensure_ascii=False)


def _read_docx(fp: str) -> str:
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    doc = _docx_mod.Document(fp)
    paragraphs = [p.text for p in doc.paragraphs]
    tables = []
    for tbl in doc.tables:
        for row in tbl.rows:
            tables.append([cell.text for cell in row.cells])
    result = "\n".join(paragraphs)
    if tables:
        result += "\n\n[Tables]\n" + json.dumps(tables, indent=2)
    return result


def _read_xlsx(fp: str) -> str:
    if not HAS_XLSX:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
    wb = openpyxl.load_workbook(fp, data_only=True)
    output: Dict[str, list] = {}
    for name in wb.sheetnames:
        ws = wb[name]
        output[name] = [list(row) for row in ws.iter_rows(values_only=True)]
    return json.dumps(output, indent=2, default=str)


def _read_pdf(fp: str) -> str:
    if not HAS_PDF_READ:
        raise RuntimeError("pdfplumber not installed. Run: pip install pdfplumber")
    pages: List[str] = []
    with pdfplumber.open(fp) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            tables = page.extract_tables() or []
            chunk = f"--- Page {i} ---\n{text}"
            if tables:
                chunk += "\n[Tables]\n" + json.dumps(tables, indent=2)
            pages.append(chunk)
    return "\n\n".join(pages)


def _read_pptx(fp: str) -> str:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = _PptxPresentation(fp)
    slides: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        slides.append(f"--- Slide {i} ---\n" + "\n".join(lines))
    return "\n\n".join(slides)


def _read_binary_b64(fp: str) -> str:
    with open(fp, "rb") as fh:
        return base64.b64encode(fh.read()).decode("utf-8")


# =============================================================================
# Format-aware WRITE Helpers
# =============================================================================

def _write_text(fp: str, content: str) -> None:
    _ensure_parent(fp)
    with open(fp, "w", encoding="utf-8") as fh:
        fh.write(content)


def _write_json(fp: str, content: str) -> None:
    parsed = json.loads(content)
    _ensure_parent(fp)
    with open(fp, "w", encoding="utf-8") as fh:
        json.dump(parsed, fh, indent=2, ensure_ascii=False)


def _write_csv(fp: str, content: str) -> None:
    rows = json.loads(content)
    _ensure_parent(fp)
    with open(fp, "w", encoding="utf-8", newline="") as fh:
        if rows and isinstance(rows[0], dict):
            writer = csv.DictWriter(fh, fieldnames=list(rows[0].keys()))
            writer.writeheader()
            writer.writerows(rows)
        else:
            csv.writer(fh).writerows(rows)


def _write_docx(fp: str, content: str) -> None:
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    _ensure_parent(fp)
    lines = content.splitlines()
    if os.path.exists(fp):
        doc = _docx_mod.Document(fp)
        for i, para in enumerate(doc.paragraphs):
            para.clear()
            if i < len(lines):
                para.add_run(lines[i])
        for line in lines[len(doc.paragraphs):]:
            doc.add_paragraph(line)
    else:
        doc = _docx_mod.Document()
        for line in lines:
            doc.add_paragraph(line)
    doc.save(fp)


def _write_xlsx(fp: str, content: str) -> None:
    if not HAS_XLSX:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
    data = json.loads(content)
    if isinstance(data, list):
        data = {"Sheet1": data}
    _ensure_parent(fp)
    wb = openpyxl.load_workbook(fp) if os.path.exists(fp) else openpyxl.Workbook()
    if not os.path.exists(fp) and "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    for sheet_name, rows in data.items():
        if sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            ws.delete_rows(1, ws.max_row)
        else:
            ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    wb.save(fp)


def _write_pdf(fp: str, content: str) -> None:
    if not HAS_PDF_WRITE:
        raise RuntimeError("reportlab not installed. Run: pip install reportlab")
    _ensure_parent(fp)
    styles = rl_styles()
    doc = SimpleDocTemplate(fp, pagesize=rl_letter)
    story = [
        RLParagraph(line if line.strip() else "&nbsp;", styles["Normal"])
        for line in content.splitlines()
    ]
    doc.build(story)


def _write_pptx(fp: str, content: str) -> None:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    _ensure_parent(fp)
    try:
        slides_data = json.loads(content)
        if not isinstance(slides_data, list):
            raise ValueError
    except (json.JSONDecodeError, ValueError):
        import re
        parts = re.split(r"---\s*Slide\s*\d+\s*---", content)
        slides_data = [{"title": "", "body": p.strip()} for p in parts if p.strip()]
    prs = _PptxPresentation(fp) if os.path.exists(fp) else _PptxPresentation()
    layout = prs.slide_layouts[1]
    for info in slides_data:
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title and isinstance(info, dict):
            slide.shapes.title.text = info.get("title", "")
        body = slide.placeholders[1]
        if body:
            body.text = info.get("body", "") if isinstance(info, dict) else str(info)
    prs.save(fp)


def _write_binary_b64(fp: str, content: str) -> None:
    _ensure_parent(fp)
    with open(fp, "wb") as fh:
        fh.write(base64.b64decode(content))


# =============================================================================
# Extension Dispatch Tables
# =============================================================================

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst",
    ".html", ".htm", ".xml", ".svg",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs",
    ".sh", ".bash", ".zsh", ".bat", ".ps1",
    ".sql", ".psql", ".r", ".rb", ".php", ".swift", ".kt", ".dart",
    ".log", ".gitignore", ".dockerignore", ".editorconfig",
}

READ_DISPATCH = {
    ".json": _read_json,
    ".csv":  _read_csv,
    ".tsv":  _read_csv,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".xls":  _read_xlsx,
    ".pdf":  _read_pdf,
    ".pptx": _read_pptx,
}

WRITE_DISPATCH = {
    ".json": _write_json,
    ".csv":  _write_csv,
    ".tsv":  _write_csv,
    ".docx": _write_docx,
    ".xlsx": _write_xlsx,
    ".xls":  _write_xlsx,
    ".pdf":  _write_pdf,
    ".pptx": _write_pptx,
}

# =============================================================================
# FILE TOOLS
# =============================================================================

async def read_file_tool(file_path: str) -> Dict[str, Any]:
    """
    Read the contents of any file inside the base directory.
    Supports: plain text, JSON, CSV/TSV, DOCX, XLSX/XLS, PDF, PPTX.
    Unknown binary files are returned as a base64-encoded string.
    Supports sub-directory paths e.g. 'subdir/report.pdf'.

    Args:
        file_path: Relative path to the file.

    Returns:
        {"status": "ok", "file_path": ..., "format": ..., "content": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(fp):
        return {"status": "error", "message": f"File not found: {file_path}"}
    if not os.path.isfile(fp):
        return {"status": "error", "message": f"Path is not a file: {file_path}"}

    ext = get_ext(fp)
    try:
        if ext in READ_DISPATCH:
            content = READ_DISPATCH[ext](fp)
            fmt = ext.lstrip(".")
        elif ext in TEXT_EXTENSIONS or ext == "":
            content = _read_text(fp)
            fmt = "text"
        else:
            try:
                content = _read_text(fp)
                fmt = "text"
            except UnicodeDecodeError:
                content = _read_binary_b64(fp)
                fmt = "base64"
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "file_path": fp, "format": fmt, "content": content}


async def write_file_tool(file_path: str, content: str) -> Dict[str, Any]:
    """
    Create or overwrite a file. Parent sub-directories are created automatically.
    Format is detected from the file extension:
      • Text/code/config files  → written as plain text.
      • .json  → content must be valid JSON.
      • .csv/.tsv  → content must be JSON: list-of-dicts or list-of-lists.
      • .docx  → plain text; each line becomes a paragraph.
      • .xlsx  → JSON: {"SheetName": [[row],...]} or [[row],...].
      • .pdf   → plain text written into a new PDF (requires reportlab).
      • .pptx  → JSON list of {"title","body"} dicts, one per slide.
      • other binary → content treated as base64 string, decoded to bytes.

    Args:
        file_path: Relative path to the file (e.g. 'reports/q1.xlsx').
        content:   Content to write (format depends on file type).

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    ext = get_ext(fp)
    try:
        if ext in WRITE_DISPATCH:
            WRITE_DISPATCH[ext](fp, content)
        elif ext in TEXT_EXTENSIONS or ext == "":
            _write_text(fp, content)
        else:
            try:
                base64.b64decode(content, validate=True)
                _write_binary_b64(fp, content)
            except Exception:
                _write_text(fp, content)
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"File written: {fp}"}


async def edit_file_tool(file_path: str, old_text: str, new_text: str) -> Dict[str, Any]:
    """
    Find every occurrence of old_text in a text-based file and replace it
    with new_text. Works on any plain-text, JSON, CSV, or code file.
    Not applicable to binary formats (DOCX, XLSX, PDF, PPTX).

    Args:
        file_path: Relative path to the file.
        old_text:  Exact substring to find.
        new_text:  Replacement string.

    Returns:
        {"status": "ok", "replacements_made": N, "message": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(fp):
        return {"status": "error", "message": f"File not found: {file_path}"}
    if not os.path.isfile(fp):
        return {"status": "error", "message": f"Not a file: {file_path}"}

    ext = get_ext(fp)
    if ext in {".docx", ".xlsx", ".xls", ".pdf", ".pptx"}:
        return {
            "status": "error",
            "message": f"edit_file_tool does not support {ext}. Use read_file_tool → modify → write_file_tool.",
        }

    try:
        original = _read_text(fp)
    except Exception as exc:
        return {"status": "error", "message": f"Cannot read file: {exc}"}

    count = original.count(old_text)
    if count == 0:
        return {"status": "error", "message": "old_text not found in file.", "replacements_made": 0}

    try:
        _write_text(fp, original.replace(old_text, new_text))
    except Exception as exc:
        return {"status": "error", "message": f"Cannot write file: {exc}"}

    return {"status": "ok", "replacements_made": count, "message": f"Replaced {count} occurrence(s)."}


async def append_file_tool(file_path: str, content: str) -> Dict[str, Any]:
    """
    Append text to the end of a file. Creates the file and any missing
    parent directories if needed. Not suitable for DOCX/XLSX/PDF/PPTX.

    Args:
        file_path: Relative path to the file.
        content:   Text to append.

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    ext = get_ext(fp)
    if ext in {".docx", ".xlsx", ".xls", ".pdf", ".pptx"}:
        return {
            "status": "error",
            "message": f"append_file_tool does not support {ext}. Use read + write instead.",
        }

    try:
        _ensure_parent(fp)
        with open(fp, "a", encoding="utf-8") as fh:
            fh.write(content)
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"Content appended to: {fp}"}


async def clear_file_tool(file_path: str) -> Dict[str, Any]:
    """
    Truncate a text-based file to zero bytes without deleting it.

    Args:
        file_path: Relative path to the file.

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(fp):
        return {"status": "error", "message": f"File not found: {file_path}"}

    ext = get_ext(fp)
    if ext in {".docx", ".xlsx", ".xls", ".pdf", ".pptx"}:
        return {
            "status": "error",
            "message": f"clear_file_tool does not support {ext}. Use write_file_tool with empty content.",
        }

    open(fp, "w").close()
    return {"status": "ok", "message": f"File cleared: {fp}"}


async def copy_file_tool(source_path: str, destination_path: str) -> Dict[str, Any]:
    """
    Copy a file from source_path to destination_path inside the base directory.
    Parent directories for the destination are created automatically.

    Args:
        source_path:      Relative path of the source file.
        destination_path: Relative path where the copy will be placed.

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        src = safe_path(source_path)
        dst = safe_path(destination_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(src):
        return {"status": "error", "message": f"Source not found: {source_path}"}
    if not os.path.isfile(src):
        return {"status": "error", "message": f"Source is not a file: {source_path}"}

    try:
        _ensure_parent(dst)
        shutil.copy2(src, dst)
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"Copied '{src}' → '{dst}'"}


async def move_file_tool(source_path: str, destination_path: str) -> Dict[str, Any]:
    """
    Move or rename a file within the base directory.
    Parent directories for the destination are created automatically.

    Args:
        source_path:      Relative path of the existing file.
        destination_path: Relative target path (new name or new location).

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        src = safe_path(source_path)
        dst = safe_path(destination_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(src):
        return {"status": "error", "message": f"Source not found: {source_path}"}
    if not os.path.isfile(src):
        return {"status": "error", "message": f"Source is not a file: {source_path}"}

    try:
        _ensure_parent(dst)
        shutil.move(src, dst)
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"Moved '{src}' → '{dst}'"}


async def delete_file_tool(file_path: str) -> Dict[str, Any]:
    """
    Permanently delete a single file from the base directory.

    Args:
        file_path: Relative path to the file.

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(fp):
        return {"status": "error", "message": f"File not found: {file_path}"}
    if not os.path.isfile(fp):
        return {"status": "error", "message": f"Not a file: {file_path}"}

    os.remove(fp)
    return {"status": "ok", "message": f"Deleted file: {fp}"}


async def file_info_tool(file_path: str) -> Dict[str, Any]:
    """
    Return metadata for a file or directory: size, timestamps,
    extension, detected format, file/directory flag.

    Args:
        file_path: Relative path to the file or directory.

    Returns:
        {"status": "ok", "file_path": ..., "extension": ..., "format": ...,
         "size_bytes": ..., "created_time": ..., "modified_time": ...,
         "is_file": ..., "is_directory": ...}
    """
    try:
        fp = safe_path(file_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(fp):
        return {"status": "error", "message": f"Path not found: {file_path}"}

    stat = os.stat(fp)
    ext  = get_ext(fp)

    if ext in READ_DISPATCH:
        fmt = ext.lstrip(".")
    elif ext in TEXT_EXTENSIONS:
        fmt = "text"
    elif os.path.isdir(fp):
        fmt = "directory"
    else:
        fmt = "binary/unknown"

    return {
        "status":        "ok",
        "file_path":     fp,
        "extension":     ext,
        "format":        fmt,
        "size_bytes":    stat.st_size,
        "created_time":  stat.st_ctime,
        "modified_time": stat.st_mtime,
        "is_file":       os.path.isfile(fp),
        "is_directory":  os.path.isdir(fp),
    }


# =============================================================================
# DIRECTORY TOOLS
# =============================================================================

async def create_directory_tool(folder_path: str) -> Dict[str, Any]:
    """
    Create a directory (including any missing intermediate sub-directories).
    Safe to call even if the directory already exists.
    Example: 'projects/2024/q1/data' creates all four levels at once.

    Args:
        folder_path: Relative path of the directory to create.

    Returns:
        {"status": "ok", "message": ..., "path": ...}
    """
    try:
        dp = safe_path(folder_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    try:
        os.makedirs(dp, exist_ok=True)
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"Directory ready: {dp}", "path": dp}


async def delete_directory_tool(folder_path: str, recursive: bool = False) -> Dict[str, Any]:
    """
    Delete a directory from the base directory.

    Args:
        folder_path: Relative path of the directory to delete.
        recursive:   If True, delete the directory and ALL its contents.
                     If False (default), directory must be empty.

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        dp = safe_path(folder_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(dp):
        return {"status": "error", "message": f"Directory not found: {folder_path}"}
    if not os.path.isdir(dp):
        return {"status": "error", "message": f"Not a directory: {folder_path}"}
    if os.path.abspath(dp) == BASE_DIR:
        return {"status": "error", "message": "Cannot delete the base (root) directory."}

    try:
        if recursive:
            shutil.rmtree(dp)
        else:
            os.rmdir(dp)
    except OSError as exc:
        return {
            "status": "error",
            "message": str(exc) + (
                "  Hint: Pass recursive=true to delete a non-empty directory."
                if not recursive else ""
            ),
        }
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"Directory deleted: {dp}"}


async def rename_directory_tool(source_path: str, destination_path: str) -> Dict[str, Any]:
    """
    Rename or move a directory within the base directory.
    Parent directories for the destination are created automatically.

    Args:
        source_path:      Relative path of the existing directory.
        destination_path: Relative path of the new location / name.

    Returns:
        {"status": "ok", "message": ...}
    """
    try:
        src = safe_path(source_path)
        dst = safe_path(destination_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(src):
        return {"status": "error", "message": f"Source not found: {source_path}"}
    if not os.path.isdir(src):
        return {"status": "error", "message": f"Source is not a directory: {source_path}"}

    try:
        _ensure_parent(dst)
        shutil.move(src, dst)
    except Exception as exc:
        return {"status": "error", "message": str(exc)}

    return {"status": "ok", "message": f"Directory moved/renamed: '{src}' → '{dst}'"}


async def list_files_tool(folder_path: str = ".") -> List[Dict[str, Any]]:
    """
    List all files (non-recursive) in a directory.

    Args:
        folder_path: Relative path of the directory to list (default: base directory root).

    Returns:
        List of {"name", "relative_path", "extension", "size_bytes"} dicts.
    """
    try:
        dp = safe_path(folder_path)
    except ValueError as e:
        return [{"status": "error", "message": str(e)}]

    if not os.path.exists(dp):
        return [{"status": "error", "message": f"Directory not found: {folder_path}"}]
    if not os.path.isdir(dp):
        return [{"status": "error", "message": f"Not a directory: {folder_path}"}]

    result = []
    for name in sorted(os.listdir(dp)):
        full = os.path.join(dp, name)
        if os.path.isfile(full):
            result.append({
                "name":          name,
                "relative_path": os.path.relpath(full, BASE_DIR),
                "extension":     get_ext(name),
                "size_bytes":    os.path.getsize(full),
            })
    return result


async def list_directories_tool(folder_path: str = ".") -> List[Dict[str, Any]]:
    """
    List all immediate sub-directories inside a directory.

    Args:
        folder_path: Relative path of the directory to inspect (default: base directory root).

    Returns:
        List of {"name", "relative_path"} dicts.
    """
    try:
        dp = safe_path(folder_path)
    except ValueError as e:
        return [{"status": "error", "message": str(e)}]

    if not os.path.exists(dp):
        return [{"status": "error", "message": f"Directory not found: {folder_path}"}]
    if not os.path.isdir(dp):
        return [{"status": "error", "message": f"Not a directory: {folder_path}"}]

    result = []
    for name in sorted(os.listdir(dp)):
        full = os.path.join(dp, name)
        if os.path.isdir(full):
            result.append({
                "name":          name,
                "relative_path": os.path.relpath(full, BASE_DIR),
            })
    return result


async def list_tree_tool(folder_path: str = ".", max_depth: int = 5) -> Dict[str, Any]:
    """
    Return a recursive tree of all files and sub-directories.
    Useful for exploring the full structure at a glance.

    Args:
        folder_path: Relative path of the root directory (default: base directory root).
        max_depth:   Levels deep to descend (default 5, max 10).

    Returns:
        {"status": "ok", "root": ..., "tree": <nested dict>}
    """
    max_depth = min(max_depth, 10)

    try:
        dp = safe_path(folder_path)
    except ValueError as e:
        return {"status": "error", "message": str(e)}

    if not os.path.exists(dp):
        return {"status": "error", "message": f"Directory not found: {folder_path}"}
    if not os.path.isdir(dp):
        return {"status": "error", "message": f"Not a directory: {folder_path}"}

    def _build(path: str, depth: int) -> Dict[str, Any]:
        name = os.path.basename(path) or path
        node: Dict[str, Any] = {"type": "directory", "name": name, "children": []}
        if depth >= max_depth:
            node["children"].append({"type": "truncated", "message": "max_depth reached"})
            return node
        try:
            entries = sorted(os.listdir(path))
        except PermissionError:
            node["children"].append({"type": "error", "message": "Permission denied"})
            return node
        for entry in entries:
            full = os.path.join(path, entry)
            if os.path.isdir(full):
                node["children"].append(_build(full, depth + 1))
            else:
                node["children"].append({
                    "type":          "file",
                    "name":          entry,
                    "relative_path": os.path.relpath(full, BASE_DIR),
                    "extension":     get_ext(entry),
                    "size_bytes":    os.path.getsize(full),
                })
        return node

    tree = _build(dp, 0)
    return {"status": "ok", "root": os.path.relpath(dp, BASE_DIR), "tree": tree}


# =============================================================================
# UTILITY TOOLS
# =============================================================================

async def supported_formats_tool() -> Dict[str, Any]:
    """
    List every supported file format, its required library,
    and whether that library is currently installed.

    Returns:
        {"status": "ok", "base_directory": ..., "read_and_write": {...}, ...}
    """
    return {
        "status":        "ok",
        "base_directory": BASE_DIR,
        "read_and_write": {
            "text/code/config": {
                "extensions": sorted(TEXT_EXTENSIONS),
                "requires":   "built-in",
                "installed":  True,
            },
            "json":     {"extensions": [".json"],         "requires": "built-in",    "installed": True},
            "csv/tsv":  {"extensions": [".csv", ".tsv"],  "requires": "built-in",    "installed": True},
            "docx":     {"extensions": [".docx"],         "requires": "python-docx", "installed": HAS_DOCX},
            "xlsx/xls": {"extensions": [".xlsx", ".xls"], "requires": "openpyxl",    "installed": HAS_XLSX},
            "pptx":     {"extensions": [".pptx"],         "requires": "python-pptx", "installed": HAS_PPTX},
        },
        "read_only": {
            "pdf": {"extensions": [".pdf"], "requires": "pdfplumber", "installed": HAS_PDF_READ},
        },
        "write_only": {
            "pdf (new file)": {"extensions": [".pdf"], "requires": "reportlab", "installed": HAS_PDF_WRITE},
        },
        "fallback": {
            "binary": {
                "extensions": ["any unrecognised extension"],
                "read":       "returned as base64 string",
                "write":      "content decoded from base64 to bytes",
                "requires":   "built-in",
            },
        },
    }


# =============================================================================
# ADK Tool Registration
# =============================================================================
adk_tools = [
    # File CRUD
    FunctionTool(read_file_tool),
    FunctionTool(write_file_tool),
    FunctionTool(edit_file_tool),
    FunctionTool(append_file_tool),
    FunctionTool(copy_file_tool),
    FunctionTool(move_file_tool),
    FunctionTool(delete_file_tool),
    FunctionTool(clear_file_tool),
    FunctionTool(file_info_tool),
    # Directory management
    FunctionTool(create_directory_tool),
    FunctionTool(delete_directory_tool),
    FunctionTool(rename_directory_tool),
    FunctionTool(list_files_tool),
    FunctionTool(list_directories_tool),
    FunctionTool(list_tree_tool),
    # Utility
    FunctionTool(supported_formats_tool),
]

# =============================================================================
# MCP Handlers
# =============================================================================

@app.list_tools()
async def list_mcp_tools() -> list[mcp_types.Tool]:
    logger.info("MCP list_tools called")
    return [adk_to_mcp_tool_type(t) for t in adk_tools]


@app.call_tool()
async def call_mcp_tool(name: str, arguments: dict) -> list[mcp_types.Content]:
    logger.info(f"MCP call_tool: '{name}'")
    tool = next((t for t in adk_tools if t.name == name), None)
    if not tool:
        return [mcp_types.TextContent(type="text", text=json.dumps({"error": f"Tool not found: {name}"}))]
    try:
        result = await tool.run_async(args=arguments, tool_context=None)
        return [mcp_types.TextContent(type="text", text=json.dumps(result, indent=2))]
    except Exception as exc:
        logger.error(f"Error executing '{name}': {exc}")
        return [mcp_types.TextContent(type="text", text=json.dumps({"error": str(exc)}))]


# =============================================================================
# STDIO Server Entry-Point
# =============================================================================

async def run_mcp_stdio_server() -> None:
    try:
        async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
            await app.run(
                read_stream,
                write_stream,
                InitializationOptions(
                    server_name=app.name,
                    server_version="0.4.0",
                    capabilities=app.get_capabilities(
                        notification_options=NotificationOptions(),
                        experimental_capabilities={},
                    ),
                ),
            )
    except (RuntimeError, KeyboardInterrupt):
        logger.info("FileSystem MCP server shutting down.")
    except BaseException as exc:
        msg = str(exc).lower()
        if "cancel scope" in msg or "generatorexit" in msg:
            logger.info("FileSystem MCP server shutting down (suppressed anyio noise).")
        else:
            raise


if __name__ == "__main__":
    try:
        asyncio.run(run_mcp_stdio_server())
    except (KeyboardInterrupt, SystemExit):
        pass
